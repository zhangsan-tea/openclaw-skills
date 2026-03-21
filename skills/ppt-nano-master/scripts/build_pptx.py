#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_pptx.py — 将一组 JPG 图片合成为 16:9 白板风 PPT

用法：
    python build_pptx.py -i slide1.jpg slide2.jpg slide3.jpg -o 输出文件.pptx
    python build_pptx.py -i slide*.jpg -o output.pptx
    python build_pptx.py --dir ./ppt_outputs/my_project/ -o output.pptx

参数：
    -i / --images   图片路径列表（按顺序排列，每张图对应一页）
    --dir           扫描目录下所有 jpg（按文件名字母序）
    -o / --output   输出 .pptx 文件路径（默认：output.pptx）
    --width         幻灯片宽度（英寸，默认 13.33，即 16:9）
    --height        幻灯片高度（英寸，默认 7.5）

依赖：
    pip install python-pptx

跨平台说明：
    - Windows / macOS / Linux 统一用 `python build_pptx.py` 调用
    - 中文路径乱码问题已在脚本内部处理，无需设置 PYTHONUTF8 环境变量
"""

import argparse
import io
import os
import sys
import glob
from pathlib import Path

# Windows 下强制 stdout/stderr 使用 UTF-8，避免中文路径乱码
if sys.platform == "win32":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")


def build_pptx(image_paths: list, output_path: str, width_in: float = 13.33, height_in: float = 7.5):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("❌ 缺少依赖，请先运行：pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    if not image_paths:
        print("❌ 没有找到任何图片，请检查路径", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    blank_layout = prs.slide_layouts[6]  # 空白版式

    for i, img_path in enumerate(image_paths):
        img_path = str(img_path)
        if not os.path.exists(img_path):
            print(f"⚠️  图片不存在，跳过：{img_path}", file=sys.stderr)
            continue
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_path,
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )
        print(f"  ✅ P{i+1}: {os.path.basename(img_path)}")

    prs.save(output_path)
    print(f"\n🎉 PPT 已生成：{output_path}")
    print(f"PPTX:{output_path}")  # 机器可读标记，供 agent 解析路径用


def main():
    parser = argparse.ArgumentParser(
        description="将 JPG 图片合成为 16:9 PPT（白板风格专用）"
    )
    parser.add_argument(
        "-i", "--images",
        nargs="+",
        help="图片路径列表（按页序排列），支持通配符展开"
    )
    parser.add_argument(
        "--dir",
        help="扫描目录下所有 .jpg 文件（按文件名排序）"
    )
    parser.add_argument(
        "-o", "--output",
        default="output.pptx",
        help="输出 .pptx 文件路径（默认：output.pptx）"
    )
    parser.add_argument(
        "--width",
        type=float,
        default=13.33,
        help="幻灯片宽度（英寸，默认 13.33）"
    )
    parser.add_argument(
        "--height",
        type=float,
        default=7.5,
        help="幻灯片高度（英寸，默认 7.5）"
    )
    args = parser.parse_args()

    # 收集图片路径
    image_paths = []

    if args.dir:
        dir_path = Path(args.dir)
        if not dir_path.is_dir():
            print(f"❌ 目录不存在：{args.dir}", file=sys.stderr)
            sys.exit(1)
        image_paths = sorted(dir_path.glob("*.jpg")) + sorted(dir_path.glob("*.jpeg"))
        print(f"📂 扫描目录 {args.dir}，找到 {len(image_paths)} 张图片")

    if args.images:
        for pattern in args.images:
            expanded = glob.glob(pattern)
            if expanded:
                image_paths.extend(sorted(expanded))
            else:
                image_paths.append(pattern)

    if not image_paths:
        parser.print_help()
        sys.exit(1)

    print(f"📊 共 {len(image_paths)} 页，开始合成...")
    build_pptx(image_paths, args.output, args.width, args.height)


if __name__ == "__main__":
    main()
