#!/usr/bin/env python3
"""
Unified document-to-markdown extraction script.
Supports: PPTX, DOCX, PDF, XLSX/XLS, images (OCR)
Outputs raw Markdown with YAML frontmatter for source tracking.
For spreadsheets, also writes one CSV per sheet alongside the MD overview.
"""

import argparse
import csv
import os
import re
import sys
from pathlib import Path


def install_deps():
    """Check and report missing dependencies with install commands."""
    deps = {
        "pptx": ("python-pptx", "from pptx import Presentation"),
        "docx": ("python-docx", "from docx import Document"),
        "pdfplumber": ("pdfplumber", "import pdfplumber"),
        "PIL": ("Pillow", "from PIL import Image"),
        "pytesseract": ("pytesseract", "import pytesseract"),
        "openpyxl": ("openpyxl", "import openpyxl"),
    }
    missing = []
    for display_name, (pkg, import_stmt) in deps.items():
        try:
            exec(import_stmt)
        except ImportError:
            missing.append(f"  pip install {pkg}  # for {display_name}")
    return missing


def extract_pptx(filepath):
    """Extract all text from a PPTX file, preserving slide structure."""
    try:
        from pptx import Presentation
    except ImportError:
        return None, "python-pptx not installed. Run: pip install python-pptx"

    prs = Presentation(filepath)
    lines = []
    lines.append(f"SLIDE_COUNT: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"{'='*60}")
        lines.append(f"SLIDE {i}")
        if slide.slide_layout:
            lines.append(f"Layout: {slide.slide_layout.name}")
        lines.append(f"{'='*60}")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(f"  {text}")
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                    lines.append(f"  | {' | '.join(cells)} |")
            if shape.shape_type == 6:  # GROUP
                for child in shape.shapes:
                    if child.has_text_frame:
                        for para in child.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                lines.append(f"    {text}")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"\n[NOTES]\n  {notes}")
        lines.append("")

    return "\n".join(lines), None


def extract_docx(filepath):
    """Extract text from a DOCX file."""
    try:
        from docx import Document
    except ImportError:
        return None, "python-docx not installed. Run: pip install python-docx"

    try:
        doc = Document(filepath)
    except Exception as e:
        return None, f"Failed to open DOCX: {e}"

    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue

        # Detect heading style
        if para.style.name.startswith("Heading"):
            level = para.style.name.split()[-1]
            try:
                level = int(level)
            except ValueError:
                level = 1
            prefix = "#" * min(level, 6)
            lines.append(f"{prefix} {text}")
        else:
            lines.append(text)

    # Tables
    for table in doc.tables:
        lines.append("")
        for row in table.rows:
            cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            lines.append(f"| {' | '.join(cells)} |")
        lines.append("")

    return "\n".join(lines), None


def extract_pdf(filepath):
    """Extract text from a PDF file (non-scanned)."""
    try:
        import pdfplumber
    except ImportError:
        return None, "pdfplumber not installed. Run: pip install pdfplumber"

    try:
        pdf = pdfplumber.open(filepath)
    except Exception as e:
        return None, f"Failed to open PDF: {e}"

    lines = []
    empty_page_count = 0

    for i, page in enumerate(pdf.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            lines.append(f"--- Page {i} ---")
            lines.append(text.strip())
            lines.append("")
        else:
            empty_page_count += 1

    pdf.close()

    result = "\n".join(lines)
    return result, (empty_page_count if empty_page_count > 0 else 0)


def _sanitize_sheet_name(name):
    """Make a sheet name safe for use as a filename component."""
    # Replace filesystem-unsafe chars with underscores
    safe = re.sub(r'[\\/:*?"<>|\s]+', '_', str(name)).strip('_')
    return safe or "sheet"


def extract_xlsx(filepath, output_dir=None):
    """
    Extract all sheets from an XLSX/XLS file.
    Writes one CSV per sheet and returns a Markdown overview with all sheets
    rendered as Markdown tables (truncated to first 200 rows per sheet for
    readability; full data lives in the CSV files).
    """
    try:
        import openpyxl
    except ImportError:
        return None, "openpyxl not installed. Run: pip install openpyxl"

    try:
        wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
    except Exception as e:
        return None, f"Failed to open spreadsheet: {e}"

    basename = os.path.splitext(os.path.basename(filepath))[0]
    csv_dir = output_dir or os.path.dirname(filepath) or "."
    os.makedirs(csv_dir, exist_ok=True)

    md_lines = []
    csv_paths = []
    sheets_with_data = 0
    total_sheets = len(wb.worksheets)

    for sheet in wb.worksheets:
        sheet_name = sheet.title
        safe_name = _sanitize_sheet_name(sheet_name)

        # Read all rows
        rows = []
        for row in sheet.iter_rows(values_only=True):
            # Convert all values to strings, replace None with empty string
            row_strs = ["" if v is None else str(v).replace("\n", " ").strip() for v in row]
            # Skip fully-empty rows
            if any(c for c in row_strs):
                rows.append(row_strs)

        if not rows:
            md_lines.append(f"## Sheet: {sheet_name}\n\n_(空 sheet，无数据)_\n")
            continue

        sheets_with_data += 1

        # Write CSV
        csv_path = os.path.join(csv_dir, f"{basename}_{safe_name}.csv")
        try:
            with open(csv_path, "w", encoding="utf-8-sig", newline="") as f:
                writer = csv.writer(f)
                writer.writerows(rows)
            csv_paths.append(csv_path)
        except Exception as e:
            md_lines.append(f"## Sheet: {sheet_name}\n\n_(CSV 写入失败: {e})_\n")
            continue

        # Build MD table (truncate for readability)
        header = rows[0]
        body = rows[1:]
        display_limit = 200
        truncated = len(body) > display_limit
        display_body = body[:display_limit]

        md_lines.append(f"## Sheet: {sheet_name}")
        md_lines.append(f"_(共 {len(body)} 行数据，CSV 完整文件: `{os.path.basename(csv_path)}`)_")
        md_lines.append("")
        md_lines.append("| " + " | ".join(header) + " |")
        md_lines.append("| " + " | ".join("---" for _ in header) + " |")
        for r in display_body:
            # Escape pipe characters inside cells
            escaped = [c.replace("|", "\\|") for c in r]
            # Pad row to header length
            while len(escaped) < len(header):
                escaped.append("")
            md_lines.append("| " + " | ".join(escaped[:len(header)]) + " |")
        if truncated:
            md_lines.append("")
            md_lines.append(f"_…后 {len(body) - display_limit} 行已省略，完整数据见 CSV 文件。_")
        md_lines.append("")

    wb.close()

    if sheets_with_data == 0:
        return None, "Spreadsheet has no data in any sheet"

    # Prepend summary
    empty_sheets = total_sheets - sheets_with_data
    summary = [
        f"# Excel 表格总览: {os.path.basename(filepath)}",
        "",
        f"- Sheet 总数: {total_sheets}",
        f"- 含数据 Sheet: {sheets_with_data}" + (f"（另 {empty_sheets} 个空 sheet 已跳过）" if empty_sheets else ""),
        f"- 生成的 CSV 文件: {len(csv_paths)}",
        "",
        "**CSV 文件清单**:",
        "",
    ]
    for p in csv_paths:
        summary.append(f"- `{os.path.basename(p)}`")
    summary.append("")
    summary.append("---")
    summary.append("")

    return "\n".join(summary) + "\n".join(md_lines), None


def extract_image(filepath):
    """OCR an image file using pytesseract."""
    try:
        from PIL import Image
    except ImportError:
        return None, "Pillow not installed. Run: pip install Pillow"
    try:
        import pytesseract
    except ImportError:
        return None, "pytesseract not installed. Run: pip install pytesseract"

    try:
        img = Image.open(filepath)
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        return text.strip(), None
    except Exception as e:
        return None, f"OCR failed: {e}"


def write_md(filepath, text, source_file, fmt_label, output_dir=None):
    """Write extracted text to a Markdown file with frontmatter header."""
    import datetime

    basename = os.path.splitext(os.path.basename(source_file))[0]
    out_name = f"{basename}_AI可读.md"

    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
        out_path = os.path.join(output_dir, out_name)
    else:
        parent = os.path.dirname(filepath) or "."
        out_path = os.path.join(parent, out_name)

    today = datetime.date.today().isoformat()
    header = f"""<!-- source: {os.path.basename(source_file)} -->
<!-- 格式: {fmt_label} -->
<!-- 转换日期: {today} -->
<!-- 转换模式: 原始文本提取（模式三） -->

"""
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(header)
        f.write(text)

    return out_path


# Mapping: extension -> (extractor, label, needs_output_dir)
EXTRACTORS = {
    ".pptx": (extract_pptx, "PPTX演示文稿", False),
    ".ppt":  (extract_pptx, "PPT演示文稿", False),
    ".docx": (extract_docx, "DOCX文档", False),
    ".doc":  (extract_docx, "DOC文档", False),
    ".pdf":  (extract_pdf, "PDF文档", False),
    ".xlsx": (extract_xlsx, "XLSX表格", True),
    ".xls":  (extract_xlsx, "XLS表格", True),
    ".png":  (extract_image, "PNG图片(OCR)", False),
    ".jpg":  (extract_image, "JPEG图片(OCR)", False),
    ".jpeg": (extract_image, "JPEG图片(OCR)", False),
    ".bmp":  (extract_image, "BMP图片(OCR)", False),
    ".tiff": (extract_image, "TIFF图片(OCR)", False),
}


def main():
    parser = argparse.ArgumentParser(
        description="Convert documents (PPTX/DOCX/PDF/XLSX/images) to Markdown"
    )
    parser.add_argument("files", nargs="*", help="Input file(s) to convert")
    parser.add_argument("-o", "--output-dir", help="Output directory for MD files")
    parser.add_argument("--check-deps", action="store_true",
                        help="Check Python dependencies only (no files needed)")
    parser.add_argument("--verbose", "-v", action="store_true",
                        help="Verbose output")

    args = parser.parse_args()

    if args.check_deps:
        missing = install_deps()
        if missing:
            print("Missing dependencies:")
            for m in missing:
                print(m)
            sys.exit(1)
        print("All dependencies installed.")
        sys.exit(0)

    if not args.files:
        parser.error("At least one input file is required (or use --check-deps)")

    results = {"success": [], "skipped": [], "failed": []}

    for filepath in args.files:
        if not os.path.exists(filepath):
            results["failed"].append((filepath, "File not found"))
            continue

        ext = os.path.splitext(filepath)[1].lower()
        if ext not in EXTRACTORS:
            results["skipped"].append((filepath, f"Unsupported format: {ext}"))
            continue

        extractor, label, needs_output_dir = EXTRACTORS[ext]

        if args.verbose:
            print(f"Extracting: {filepath} ({label})...")

        if needs_output_dir:
            text, error_or_info = extractor(filepath, args.output_dir)
        else:
            text, error_or_info = extractor(filepath)

        if text is None:
            results["failed"].append((filepath, error_or_info or "Unknown error"))
            continue

        # Check for empty/minimal output
        if not text or len(text.strip()) < 20:
            reason = f"Empty or minimal content ({len(text.strip())} chars)"
            if isinstance(error_or_info, int) and error_or_info > 0:
                if ext == ".pdf":
                    # Check if it's a scanned PDF
                    pct = error_or_info / max(1, error_or_info + 1)
                    reason = f"Likely scanned PDF: {error_or_info} empty pages"
            results["skipped"].append((filepath, reason))
            continue

        out_path = write_md(filepath, text, filepath, label, args.output_dir)
        results["success"].append((filepath, out_path))

        if args.verbose:
            size_kb = os.path.getsize(out_path) / 1024
            print(f"  -> {out_path} ({size_kb:.1f} KB)")

    # Summary
    print(f"\n{'='*50}")
    print(f"Summary: {len(results['success'])} success, "
          f"{len(results['skipped'])} skipped, {len(results['failed'])} failed")

    if results["skipped"]:
        print(f"\nSkipped ({len(results['skipped'])}):")
        for f, reason in results["skipped"]:
            print(f"  {os.path.basename(f)}: {reason}")

    if results["failed"]:
        print(f"\nFailed ({len(results['failed'])}):")
        for f, reason in results["failed"]:
            print(f"  {os.path.basename(f)}: {reason}")

    sys.exit(0 if not results["failed"] else 1)


if __name__ == "__main__":
    main()
